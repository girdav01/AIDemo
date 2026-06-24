#!/usr/bin/env python3
"""Generate the booth big-screen deck (looping slides) as an editable .pptx.

Mirrors the booth runbook / screen content: eight challenges along the
Visibility -> Control -> Governance spine, persona routing (Execs / CIO / CAIO /
CISO, AI Builders, Walk-ups), and the "one per layer" passport completion.

Run:  python tools/generate_screen_deck.py
      -> dist/AI4_VisionOne_Challenge_BoothScreen.pptx

Why a generator: the original deck was a binary/legacy PptxGenJS export that
can't be text-edited in place. This rebuilds an editable, regenerable version in
the booth theme. Tweak the data tables below and re-run.
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# Booth theme
INK = RGBColor(0x0D, 0x11, 0x17)
PANEL = RGBColor(0x16, 0x1B, 0x22)
RED = RGBColor(0xD7, 0x19, 0x20)
TEXT = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x93, 0xA1, 0xB1)
GOLD = RGBColor(0xE3, 0xB3, 0x41)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EVENT = "Trend Vision One™  ·  Ai4 2026  ·  The Venetian, Las Vegas"

CHALLENGES = [
    ("1  Break the Bot", "Control", "Jailbreak a live chatbot to leak a planted secret", "AI Guard blocks you and logs the attempt"),
    ("2  Stop the Leak", "Control", "Coax the app into spilling fake PII or source code", "AI Guard redacts the sensitive data in real time"),
    ("3  Find the Flaw", "Visibility", "Run AI Scanner on a vulnerable model pre-deploy", "You name the OWASP LLM risk behind the top finding"),
    ("4  Trace the Poison", "Visibility", "Catch a hardcoded secret + bad dependency", "You trace it to a downstream app via the SBOM"),
    ("5  Shadow AI Hunt", "Control", "Spot unsanctioned GenAI use, then set a policy", "Your policy blocks the next risky prompt"),
    ("6  Tame the Agent", "Governance", "Push a rogue agent toward an unauthorized action", "Governance denies it — with a full audit trail"),
    ("7  Watch the MCP Wire", "Governance", "Spot a rogue MCP tool-call at the gateway", "The call is blocked at the gateway and logged"),
    ("8  Boss Level", "Capstone", "Speed-run scan → protect → validate → improve", "All four steps done before the timer ends"),
]

PERSONAS = [
    ("Execs / CIO / CAIO / CISO", "Start at Stop the Leak (#2) and Shadow AI Hunt (#5) — fast, business-relevant wins."),
    ("AI Builders", "Find the Flaw (#3), Trace the Poison (#4), Tame the Agent (#6), Watch the MCP Wire (#7), Boss Level (#8)."),
    ("Everyone / walk-ups", "Break the Bot (#1) is the magnet — lead with it whenever the booth is quiet."),
]


def _bg(slide, color=INK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _line(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, first=False, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Arial"
    return p


def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    tf = _box(s, 0.8, 2.2, 11.7, 3)
    _line(tf, "AI Security Challenge", 60, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    _line(tf, EVENT, 22, MUTED, align=PP_ALIGN.CENTER, space_after=0)
    tf2 = _box(s, 0.8, 4.4, 11.7, 1)
    _line(tf2, "Can you break the bot?", 28, RED, bold=True, align=PP_ALIGN.CENTER, first=True)


def statement_slide(prs, line1, line2, sub=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    tf = _box(s, 0.8, 2.6, 11.7, 2.5)
    _line(tf, line1, 54, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    _line(tf, line2, 54, RED, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        _line(tf, sub, 22, MUTED, align=PP_ALIGN.CENTER)


def bullets_slide(prs, heading, bullets, footer=None, footer_color=GOLD):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    th = _box(s, 0.8, 0.5, 11.7, 1)
    _line(th, heading, 36, RED, bold=True, first=True)
    tf = _box(s, 1.0, 1.7, 11.3, 5)
    for i, b in enumerate(bullets):
        _line(tf, "•  " + b, 24, TEXT, first=(i == 0), space_after=12)
    if footer:
        tff = _box(s, 0.8, 6.6, 11.7, 0.8)
        _line(tff, footer, 22, footer_color, bold=True, align=PP_ALIGN.CENTER, first=True)


def _cell(cell, text, size, color, bold=False, fill=None, align=PP_ALIGN.LEFT):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Arial"


def challenges_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    th = _box(s, 0.5, 0.3, 12.3, 0.9)
    _line(th, "Eight Challenges · Visibility → Control → Governance", 30, RED, bold=True, first=True)

    rows, cols = len(CHALLENGES) + 1, 4
    gshape = s.shapes.add_table(rows, cols, Inches(0.4), Inches(1.15), Inches(12.5), Inches(5.9))
    table = gshape.table
    table.columns[0].width = Inches(2.7)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(4.3)
    table.columns[3].width = Inches(4.0)
    headers = ["Challenge", "Layer", "Your mission in Vision One", "You clear it when…"]
    for c, h in enumerate(headers):
        _cell(table.cell(0, c), h, 14, WHITE, bold=True, fill=RED)
    for i, (name, layer, mission, clears) in enumerate(CHALLENGES, start=1):
        _cell(table.cell(i, 0), name, 13, TEXT, bold=True)
        _cell(table.cell(i, 1), layer, 12, GOLD, bold=True)
        _cell(table.cell(i, 2), mission, 12, TEXT)
        _cell(table.cell(i, 3), clears, 12, MUTED)


def personas_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    th = _box(s, 0.8, 0.5, 11.7, 1)
    _line(th, "Find your track", 36, RED, bold=True, first=True)
    tf = _box(s, 1.0, 1.8, 11.3, 4.5)
    first = True
    for who, where in PERSONAS:
        p = _line(tf, who, 26, GOLD, bold=True, first=first, space_after=2)
        first = False
        _line(tf, where, 20, TEXT, space_after=18)
    tff = _box(s, 0.8, 6.5, 11.7, 0.8)
    _line(tff, "Everyone hits one Visibility, one Control, and one Governance station.",
          22, MUTED, align=PP_ALIGN.CENTER, first=True)


def contact_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    tf = _box(s, 0.8, 2.6, 11.7, 2.5)
    _line(tf, "David Girard", 40, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    _line(tf, "david_girard@trendmicro.com", 24, MUTED, align=PP_ALIGN.CENTER)
    _line(tf, "AI Fearlessly — clear a challenge, then talk to our team.",
          24, RED, bold=True, align=PP_ALIGN.CENTER)


def build(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    statement_slide(prs, "Can you", "break the bot?")
    bullets_slide(
        prs, "How to Play",
        [
            "Grab a Challenge Passport at the counter.",
            "Clear any station in about 5 minutes with a TrendAI guide.",
            "Earn a stamp in each layer — Visibility · Control · Governance.",
            "Watch your name climb the live leaderboard on this screen.",
            "One challenge per layer fills your passport and enters the grand-prize draw.",
        ],
        footer="Only 37% of organizations test their AI before they ship it. Be the exception.",
    )
    challenges_slide(prs)
    personas_slide(prs)
    statement_slide(prs, "Don't just build AI.", "Secure it.")
    bullets_slide(
        prs, "Prizes & Leaderboard",
        [
            "Beat one challenge → instant swag.",
            "Full passport (one per layer) → premium swag + grand-prize draw.",
            "Daily leaderboard top 3 → headline prize, drawn at 4 PM each day.",
            "Boss Level finishers → exclusive “AI Fearlessly” challenge coin.",
        ],
        footer="Don't just build AI. Secure it. — AI Fearlessly",
        footer_color=RED,
    )
    contact_slide(prs)

    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    print("Wrote", path, "(", len(prs.slides._sldIdLst), "slides )")


if __name__ == "__main__":
    build("dist/AI4_VisionOne_Challenge_BoothScreen.pptx")
